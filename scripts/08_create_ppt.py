"""
나침반 (NACHIMBAN) - 대회 제출용 PPT 생성 (세련된 디자인)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

OUTPUT_PATH = Path("docs/나침반_NACHIMBAN_제출자료.pptx")
CHART_DIR = Path("docs/charts")

SLIDE_WIDTH = Cm(33.867)  # 16:9 와이드
SLIDE_HEIGHT = Cm(19.05)

# 컬러 팔레트
PRIMARY = RGBColor(0x1D, 0x4E, 0xD8)    # 진한 파란
ACCENT = RGBColor(0x3B, 0x82, 0xF6)     # 밝은 파란
DARK = RGBColor(0x0F, 0x17, 0x2A)       # 거의 검정
TEXT = RGBColor(0x1E, 0x29, 0x3B)       # 본문
GRAY = RGBColor(0x64, 0x74, 0x8B)       # 부제
LIGHT = RGBColor(0x94, 0xA3, 0xB8)      # 연한 회색
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_BLUE = RGBColor(0xEF, 0xF6, 0xFF)    # 연한 파란 배경
GREEN = RGBColor(0x05, 0x96, 0x69)      # 초록


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_left_accent(slide):
    """왼쪽에 파란 세로 바 추가"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(0.8), SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()


def add_page_number(slide, num):
    """우하단 페이지 번호"""
    txBox = slide.shapes.add_textbox(Cm(31), Cm(17.5), Cm(2), Cm(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT
    p.alignment = PP_ALIGN.RIGHT


def make_title_slide(prs):
    """표지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    
    # 왼쪽 파란 원형 데코
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(-5), Cm(-3), Cm(18), Cm(18))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    
    # 서비스명
    txBox = slide.shapes.add_textbox(Cm(10), Cm(4), Cm(22), Cm(4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "나침반"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "NACHIMBAN"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT
    p2.space_before = Pt(4)
    
    # 부제
    txBox2 = slide.shapes.add_textbox(Cm(10), Cm(10), Cm(22), Cm(3))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = "AI 기반 교육 공공데이터 활용"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    
    p4 = tf2.add_paragraph()
    p4.text = "고등학교 진로 매칭 서비스"
    p4.font.size = Pt(16)
    p4.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    
    # 하단 URL
    txBox3 = slide.shapes.add_textbox(Cm(10), Cm(15), Cm(22), Cm(2))
    tf3 = txBox3.text_frame
    p5 = tf3.paragraphs[0]
    p5.text = "https://nachimban-nine.vercel.app"
    p5.font.size = Pt(12)
    p5.font.color.rgb = ACCENT
    
    p6 = tf3.add_paragraph()
    p6.text = "교육공공데이터 AI 활용 대회 | AI 활용 아이디어 기획"
    p6.font.size = Pt(11)
    p6.font.color.rgb = LIGHT


def make_content_slide(prs, page_num, title, content_lines):
    """본문 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_left_accent(slide)
    add_page_number(slide, page_num)
    
    # 제목
    txBox = slide.shapes.add_textbox(Cm(2.5), Cm(1), Cm(28), Cm(2.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK
    
    # 제목 아래 라인
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2.5), Cm(3.2), Cm(4), Cm(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    # 본문
    txBox2 = slide.shapes.add_textbox(Cm(2.5), Cm(4), Cm(29), Cm(14))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    for text in content_lines:
        p = tf2.add_paragraph()
        p.space_after = Pt(3)
        
        if text == "":
            p.space_after = Pt(8)
            continue
        elif text.startswith("▶"):
            p.text = text
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = PRIMARY
            p.space_before = Pt(12)
        elif text.startswith("  •"):
            p.text = text
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT
        elif text.startswith("■"):
            p.text = text
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = DARK
            p.space_before = Pt(8)
        else:
            p.text = text
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT
    
    return slide


def make_image_slide(prs, page_num, title, image_path, caption):
    """이미지 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_left_accent(slide)
    add_page_number(slide, page_num)
    
    # 제목
    txBox = slide.shapes.add_textbox(Cm(2.5), Cm(0.8), Cm(28), Cm(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK
    
    # 이미지
    img_path = Path(image_path)
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Cm(4), Cm(3), Cm(26), Cm(13.5))
    
    # 캡션
    txBox2 = slide.shapes.add_textbox(Cm(2.5), Cm(17), Cm(29), Cm(1.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = caption
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.alignment = PP_ALIGN.CENTER
    
    return slide


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # 1. 표지
    make_title_slide(prs)
    
    # 2. 활용 데이터 (규정: 첫 페이지)
    make_content_slide(prs, 1, "활용 교육 공공데이터", [
        "■ 데이터 1: 전국 고등학교 기본정보 (2,407개 학교)",
        "  • 출처: 나이스(NEIS) 교육정보 개방 포털",
        "  • URL: https://open.neis.go.kr/hub/schoolInfo",
        "  • 라이선스: 공공누리 제1유형 (출처표시)",
        "  • 내용: 학교명, 유형, 주소, 홈페이지, 설립구분, 남녀공학구분",
        "",
        "■ 데이터 2: 고등학교별 개설 교과목 (2,307개 학교, 161,317건)",
        "  • 출처: 나이스(NEIS) 교육정보 개방 포털",
        "  • URL: https://open.neis.go.kr/hub/hisTimetable",
        "  • 라이선스: 공공누리 제1유형 (출처표시)",
        "  • 내용: 학교별 실제 개설 교과목명 (2025년 4월 기준 시간표)",
        "",
        "■ 데이터 3: 직업 상세정보 + 능력/지식 점수 (462개 직업)",
        "  • 출처: 워크넷 work24 OpenAPI",
        "  • URL: https://www.work24.go.kr/cm/openApi/call/wk/callOpenApiSvcInfo212L01.do",
        "  • 라이선스: 공공누리 제1유형 (출처표시)",
        "  • 내용: 직업명, 개요, 임금, 전망, 능력 35개 + 지식 30개 항목 점수",
    ])
    
    # 3. 문제 정의
    make_content_slide(prs, 2, "문제 정의", [
        "▶ 현재 진로 교육의 문제점",
        "",
        "  • 전국 2,407개 고등학교 정보가 여러 기관에 산재 → 통합 비교 불가",
        "  • 기존 MBTI 검사: 16유형으로만 분류 → 경계값 학생의 개인차 미반영",
        "  • 진로 선택이 데이터 기반이 아닌 주변 평판·막연한 인식에 의존",
        "  • 진로상담 인프라 부족 지역(농어촌 등)의 정보 접근성 격차",
        "",
        "▶ 해결 목표",
        "",
        "  • 교육 공공데이터를 AI로 통합 분석하여 맞춤형 진로 매칭 제공",
        "  • MBTI 연속 스펙트럼(4축 0~100)으로 256가지+ 개인차 반영",
        "  • 웹 기반 무료 서비스 → 지역·소득 무관 동등한 정보 접근",
    ])
    
    # 4. 서비스 개요
    make_content_slide(prs, 3, "서비스 개요: 나침반(NACHIMBAN)", [
        "▶ 핵심 컨셉: Holland + MBTI 연속 스펙트럼 + 공공데이터 AI 매칭",
        "",
        "■ MBTI 연속 스펙트럼 (기존 대비 차별점)",
        "  • 기존: E or I (이분법) → 나침반: EI축 0~100 연속값",
        "  • 4축 × 4단계 = 256가지 조합 (실제로는 무한한 연속값)",
        "  • 예: EI=35(약한E), SN=80(강한N), TF=25(강한T), JP=60(약한P)",
        "",
        "■ 이중 매칭 구조",
        "  • 1차: Holland(RIASEC) 코사인 유사도 → 직업군 필터 (가중치 60%)",
        "  • 2차: MBTI 유클리드 거리 → 직업 내 세부 매칭 (가중치 40%)",
        "  • 3차: 공공데이터 기반 학과 → 교과목 → 학교 연결",
        "",
        "■ 결과 제공",
        "  • 추천 진로 Top 5 + 목표 매칭도 비교 + 맞춤 학교 + 교육과정 로드맵",
    ])
    
    # 5. AI 활용 과정
    make_content_slide(prs, 4, "AI 활용 과정", [
        "▶ Step 1: 데이터 수집 및 전처리",
        "  • 나이스 API → 고등학교 기본정보 + 시간표(개설 교과목) 수집",
        "  • 워크넷 API → 462개 직업별 능력(35개)/지식(30개) 점수 수집",
        "",
        "▶ Step 2: AI 모델링 (공공데이터 → Holland/MBTI 프로필 산출)",
        "  • Holland: 워크넷 능력 키워드를 RIASEC 6유형에 매핑 → 유형별 평균 점수",
        "  • MBTI: 양방향 능력 키워드 비율로 4축 연속값 산출",
        "  • 공식: 축 값 = 두번째방향평균 / (첫번째 + 두번째) × 100",
        "",
        "▶ Step 3: 추론 및 매칭",
        "  • 학생 프로필 ↔ 직업 프로필 유사도 계산 → 적합 직업 순위화",
        "  • 적합 직업 → 관련 학과 → 필수 교과목 → 해당 과목 개설 학교 필터",
        "  • 성별 필터 + 거주지 근접도(동/구 단위) 반영 → 최종 학교 추천",
    ])
    
    # 6. 사용자 플로우
    make_content_slide(prs, 5, "서비스 사용자 플로우", [
        "▶ 입력 단계 (약 8~10분 소요)",
        "  • [Step 1] 목표 입력 — 희망 학교/학과/직업 자동완성 선택 (건너뛰기 가능)",
        "  • [Step 2] Holland 흥미 검사 — 18문항 · 5점 척도 → RIASEC 6축 점수",
        "  • [Step 3] MBTI 스펙트럼 검사 — 32문항 · 7점 척도 → 4축 연속값(0~100)",
        "  • [Step 4] 추가 정보 — 성별, 거주지(시도/시군구/동), 학년",
        "",
        "▶ 결과 대시보드 (AI 분석 후 즉시 제공)",
        "  • 📋 종합 의견 — 핵심 분석 결과 한줄 요약",
        "  • 📊 프로필 시각화 — Holland 레이더 차트 + MBTI 스펙트럼 바(수치 포함)",
        "  • 🎯 목표 매칭도 — 입력 목표 vs AI 추천 적합도 비교 (0~100%)",
        "  • 💼 추천 진로 Top 5 — 직업 설명 + 적합도 + 추천 이유 + 관련 학과",
        "  • 🏫 추천 학교 3~10개 — 유형별 차등 정보, 더보기 지원",
        "  • 📚 교육과정 로드맵 — 고1→고2→고3 추천 교과 + 비교과",
    ])
    
    # 7~11. 데이터 분석 차트
    make_image_slide(prs, 6, "데이터 분석 ① 전국 고등학교 유형별 분포",
        CHART_DIR / "01_school_types.png",
        "일반고 1,645개(68%) · 특성화고 490개(20%) · 특목고 160개(7%) · 자율고 112개(5%) | 나이스 공공데이터 기반")
    
    make_image_slide(prs, 7, "데이터 분석 ② 시도별 고등학교 수",
        CHART_DIR / "02_schools_by_region.png",
        "경기도(505개) > 서울(319개) > 경남(195개) | 수도권 집중 및 지역 간 인프라 격차 확인")
    
    make_image_slide(prs, 8, "데이터 분석 ③ Holland 유형별 직업 분포",
        CHART_DIR / "04_holland_distribution.png",
        "현실형(R) 37%로 최다 | 워크넷 462개 직업의 능력/지식 데이터 기반 Holland 코드 산출 결과")
    
    make_image_slide(prs, 9, "데이터 분석 ④ 직업 대분류별 MBTI 스펙트럼",
        CHART_DIR / "05_mbti_by_category.png",
        "연구직: I+T 성향 | 교육직: E+F 성향 | 예술직: N+P 성향 | 대분류별 MBTI 패턴이 명확히 구분됨")
    
    make_image_slide(prs, 10, "데이터 분석 ⑤ 전국 고등학교 인기 개설 과목",
        CHART_DIR / "07_popular_subjects.png",
        "한국사·영어Ⅰ·수학Ⅰ = 대부분 개설 | 물리학Ⅱ·공학일반 등 진로 특화 과목은 소수 학교만 개설 → 학교 선택의 근거")
    
    # 12. 기대효과
    make_content_slide(prs, 11, "기대효과 및 사회적 가치", [
        "▶ 기대효과",
        "  • 학생: 데이터 기반 객관적 진로 탐색, 개인 특성에 맞는 학교 발견",
        "  • 학부모: 자녀 적성 데이터로 교육 의사결정 근거 확보",
        "  • 교사: 진로지도 보조 자료, 학생별 맞춤 상담 근거",
        "  • 정책: 지역별 교과목 개설 격차 데이터 기반 파악",
        "",
        "▶ 사회적 가치",
        "  • 교육 형평성 — 인프라 부족 지역에도 동일 품질의 AI 추천 제공",
        "  • 정보 접근성 — 나이스+워크넷 공공데이터를 하나의 서비스로 통합",
        "  • 개인 맞춤화 — 256가지+ MBTI 조합으로 획일적 진로지도 탈피",
        "  • 공공데이터 활용 — 세금으로 구축된 데이터의 교육 현장 환원",
        "  • 의사결정 지원 — 감이 아닌 데이터 기반 진로 선택 문화 조성",
    ])
    
    # 13. 기술 아키텍처
    make_content_slide(prs, 12, "기술 아키텍처", [
        "▶ 시스템 구성",
        "  • Frontend: React + TypeScript + Tailwind CSS + Recharts → Vercel 배포",
        "  • Backend: Python FastAPI + AI 매칭 엔진 → Render 배포",
        "  • 데이터: 나이스 API + 워크넷 API → CSV 가공 → 매칭 엔진에 로드",
        "",
        "▶ AI 매칭 엔진 핵심 로직",
        "  • Holland: 학생 6축 벡터 ↔ 직업 6축 벡터 → 코사인 유사도(0~100)",
        "  • MBTI: 학생 4축 ↔ 직업 4축 → 유클리드 거리 → 유사도(0~100)",
        "  • 최종 점수 = Holland × 0.6 + MBTI × 0.4",
        "",
        "▶ 학교 매칭 점수",
        "  • 기본: 진로 필수 교과목 개설 비율 (개설 수 / 전체 필수 과목 수) × 100",
        "  • 보너스: 특목고 +5점, 자율고 +3점, 같은 동 +8점, 같은 구 +3점",
        "  • 필터: 성별(남→여학교 제외), 거주지(시군구 우선 → 시도 확장)",
    ])
    
    # 14. 향후 계획
    make_content_slide(prs, 13, "향후 발전 방향", [
        "▶ 데이터 확장",
        "  • 커리어넷 API 연동 (승인 대기 중) → 학과-직업 매핑 정확도 향상",
        "  • 졸업생 진로 현황 데이터 → 학교별 실제 진학 실적 기반 추천",
        "  • 대학 입시 전형 정보 연계 → 학생부교과/종합 가이드",
        "",
        "▶ 기능 고도화",
        "  • 사용자 피드백 수집 → 추천 모델 지속 개선 (A/B 테스트)",
        "  • 학부모/교사용 대시보드 → 교육 현장 활용도 제고",
        "  • 모바일 반응형 → 스마트폰 접근성 확대",
        "",
        "▶ 서비스 확장 로드맵",
        "  • Phase 1 (현재): 중학생 → 고등학교 매칭",
        "  • Phase 2: 고등학생 → 대학/학과 매칭",
        "  • Phase 3: 대학생 → 취업/진로 매칭",
    ])
    
    # 15. 부록: 생성형 AI 활용
    make_content_slide(prs, "", "부록: 생성형 AI 활용 정보", [
        "▶ 활용 도구",
        "  • Claude (Anthropic) — Kiro IDE 환경에서 대화형으로 활용",
        "",
        "▶ 활용 범위",
        "  • 프로젝트 기획 및 아키텍처 설계 지원",
        "  • 공공데이터 API 수집 스크립트 작성 (Python)",
        "  • Holland/MBTI 매핑 알고리즘 설계 및 구현",
        "  • Frontend(React) / Backend(FastAPI) 코드 작성",
        "  • EDA 시각화 스크립트 작성",
        "  • 보고서 및 발표자료 초안 작성",
        "",
        "▶ 산출물",
        "  • 전체 소스코드: https://github.com/soobeen0822/nachimban",
        "  • 라이브 서비스: https://nachimban-nine.vercel.app",
        "  • API 서버: https://nachimban-api.onrender.com",
    ])
    
    # 저장
    prs.save(str(OUTPUT_PATH))
    print(f"✅ PPT 생성 완료: {OUTPUT_PATH}")
    print(f"   총 {len(prs.slides)}페이지 (표지 포함, 부록 제외 본문 13매)")


if __name__ == "__main__":
    main()
